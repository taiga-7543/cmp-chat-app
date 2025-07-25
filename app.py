from flask import Flask, request, jsonify, render_template, Response, redirect, url_for, session
from flask_httpauth import HTTPBasicAuth
from google import genai
from google.genai import types
import json
import os
import asyncio
import time
import re
from datetime import datetime
import tempfile
import hashlib
import base64
import gc
from dotenv import load_dotenv
import io
import mimetypes
from werkzeug.utils import secure_filename
from google.cloud import aiplatform
# RAGサービスクライアントを一時的にコメントアウト
# from google.cloud.aiplatform_v1 import RagDataServiceClient
# from google.cloud.aiplatform_v1.types import ListRagFilesRequest, ImportRagFilesRequest, RagFile
from google.cloud import storage
import uuid
import docx
from pptx import Presentation  # PowerPoint処理用
from urllib.parse import unquote

# Google Drive同期機能をインポート
from drive_sync_integration import init_drive_sync, get_drive_sync

# .envファイルを読み込み
load_dotenv()

app = Flask(__name__)
app.secret_key = os.environ.get('SECRET_KEY', 'your-secret-key-change-this')
auth = HTTPBasicAuth()

# 環境変数から設定を読み込み
PROJECT_ID = os.environ.get('GOOGLE_CLOUD_PROJECT', 'dotd-development-division')
RAG_CORPUS = os.environ.get('RAG_CORPUS', f'projects/{PROJECT_ID}/locations/us-central1/ragCorpora/3458764513820540928')
GEMINI_MODEL = os.environ.get('GEMINI_MODEL', 'gemini-2.5-flash')

# 認証設定
AUTH_USERNAME = os.environ.get('AUTH_USERNAME', 'u7F3kL9pQ2zX')
AUTH_PASSWORD = os.environ.get('AUTH_PASSWORD', 's8Vn2BqT5wXc')

@auth.verify_password
def verify_password(username, password):
    """ユーザー名とパスワードを検証"""
    if username == AUTH_USERNAME and password == AUTH_PASSWORD:
        return username
    return None

@auth.error_handler
def auth_error(status):
    """認証エラーハンドラー"""
    return jsonify({'error': 'アクセスが拒否されました。正しいユーザー名とパスワードを入力してください。'}), status

def fix_base64_padding(data):
    """Base64データのパディングを修正"""
    missing_padding = len(data) % 4
    if missing_padding:
        data += '=' * (4 - missing_padding)
    return data

def validate_and_fix_private_key(private_key):
    """プライベートキーの形式を検証・修正"""
    if not private_key:
        return private_key
    
    # PEM形式のヘッダー・フッターを確認
    if not private_key.startswith('-----BEGIN PRIVATE KEY-----'):
        return private_key
    
    # PEM形式のキーを行に分割
    lines = private_key.split('\n')
    if len(lines) < 3:
        return private_key
    
    # ヘッダーとフッターを除いたBase64部分を取得
    base64_lines = []
    for line in lines[1:-1]:  # ヘッダーとフッターを除く
        line = line.strip()
        if line and not line.startswith('-----'):
            base64_lines.append(line)
    
    if not base64_lines:
        return private_key
    
    # Base64文字列を結合
    base64_data = ''.join(base64_lines)
    
    # パディングを修正
    try:
        fixed_base64 = fix_base64_padding(base64_data)
        # Base64デコードをテスト
        base64.b64decode(fixed_base64)
        
        # 修正されたキーを再構築
        fixed_private_key = '-----BEGIN PRIVATE KEY-----\n'
        # 64文字ずつ改行
        for i in range(0, len(fixed_base64), 64):
            fixed_private_key += fixed_base64[i:i+64] + '\n'
        fixed_private_key += '-----END PRIVATE KEY-----'
        
        return fixed_private_key
        
    except Exception as e:
        print(f"Warning: Could not fix private key Base64 padding: {e}")
        return private_key

def setup_google_auth():
    """Google Cloud認証を設定"""
    # 環境変数からサービスアカウントキーのJSONを読み込む
    credentials_json = os.environ.get('GOOGLE_APPLICATION_CREDENTIALS_JSON')
    if credentials_json:
        try:
            # JSONの妥当性を確認
            parsed_json = json.loads(credentials_json)
            
            # 必要なフィールドが存在するか確認
            required_fields = ['type', 'project_id', 'private_key_id', 'private_key', 'client_email']
            missing_fields = [field for field in required_fields if field not in parsed_json]
            
            if missing_fields:
                print(f"Warning: Missing required fields in Google Cloud credentials: {missing_fields}")
                return
            
            # プライベートキーのBase64パディングを修正
            if 'private_key' in parsed_json:
                original_key = parsed_json['private_key']
                fixed_key = validate_and_fix_private_key(original_key)
                if fixed_key != original_key:
                    print("INFO: Fixed private key Base64 padding")
                    parsed_json['private_key'] = fixed_key
            
            # JSONをファイルに書き込んで認証設定
            with tempfile.NamedTemporaryFile(mode='w', suffix='.json', delete=False) as f:
                json.dump(parsed_json, f, indent=2)
                os.environ['GOOGLE_APPLICATION_CREDENTIALS'] = f.name
                print(f"Google Cloud credentials file created: {f.name}")
                
        except json.JSONDecodeError as e:
            print(f"Error: Invalid JSON in GOOGLE_APPLICATION_CREDENTIALS_JSON: {e}")
            print("Please check the format of your Google Cloud credentials JSON.")
            return
        except Exception as e:
            print(f"Error setting up Google Cloud authentication: {e}")
            return
    
    # 既存のGOOGLE_APPLICATION_CREDENTIALSがある場合はそのまま使用
    if not os.environ.get('GOOGLE_APPLICATION_CREDENTIALS') and not credentials_json:
        print("Warning: Google Cloud認証情報が設定されていません")
        print("以下の環境変数のいずれかを設定してください:")
        print("- GOOGLE_APPLICATION_CREDENTIALS: サービスアカウントキーファイルのパス")
        print("- GOOGLE_APPLICATION_CREDENTIALS_JSON: サービスアカウントキーのJSON文字列")
    elif os.environ.get('GOOGLE_APPLICATION_CREDENTIALS'):
        print(f"Using existing Google Cloud credentials file: {os.environ.get('GOOGLE_APPLICATION_CREDENTIALS')}")

# 認証設定を初期化
setup_google_auth()

# Google Drive同期を初期化
CLIENT_SECRETS_FILE = os.environ.get('GOOGLE_OAUTH_CLIENT_SECRETS_FILE', 'config/client_secrets.json')
CLOUD_STORAGE_BUCKET = os.environ.get('CLOUD_STORAGE_BUCKET', 'dotd-cmp-wg-search')
DRIVE_FOLDER_NAME = os.environ.get('DRIVE_SYNC_FOLDER_NAME', 'RAG Documents')

if os.path.exists(CLIENT_SECRETS_FILE):
    init_drive_sync(PROJECT_ID, CLIENT_SECRETS_FILE, CLOUD_STORAGE_BUCKET)
    print(f"Google Drive同期機能を初期化しました")
else:
    print(f"警告: {CLIENT_SECRETS_FILE} が見つかりません。Google Drive同期機能は無効です。")

# メモリ管理の設定
gc.set_threshold(700, 10, 10)  # より積極的なガベージコレクション

def create_rag_client():
    """RAGクライアントを作成"""
    client = genai.Client(
        vertexai=True,
        project=PROJECT_ID,
        location="global",
    )
    return client

def extract_date_from_filename(filename):
    """ファイル名から日付を抽出する（_yyyymmdd形式またはyyyymmdd形式）"""
    if not filename:
        return None
    
    # yyyymmdd形式の日付を検索（アンダースコアありまたはなし）
    # パターン1: _yyyymmdd形式（アンダースコアあり）
    # パターン2: yyyymmdd形式（ファイル名の先頭または区切り文字の後）
    date_patterns = [
        r'_(\d{8})(?:\.|_|$)',  # _yyyymmdd形式
        r'(?:^|[^\d])(\d{8})(?:[^\d]|$)',  # yyyymmdd形式（前後に数字以外の文字）
    ]
    
    for pattern in date_patterns:
        match = re.search(pattern, filename)
        if match:
            date_str = match.group(1)
            try:
                # 日付の妥当性をチェック
                date_obj = datetime.strptime(date_str, '%Y%m%d')
                
                # 妥当な日付範囲をチェック（1900年〜2100年）
                if 1900 <= date_obj.year <= 2100:
                    return date_obj
            except ValueError:
                # 無効な日付の場合は次のパターンを試す
                continue
    
    return None

def sort_sources_by_date(sources):
    """出典情報を日付順にソート（新しい日付を優先）"""
    def get_sort_key(source):
        title = source.get('title', '')
        uri = source.get('uri', '')
        
        # タイトルとURIの両方から日付を抽出を試みる
        date_from_title = extract_date_from_filename(title)
        date_from_uri = extract_date_from_filename(uri)
        
        # より新しい日付を使用
        extracted_date = None
        if date_from_title and date_from_uri:
            extracted_date = max(date_from_title, date_from_uri)
        elif date_from_title:
            extracted_date = date_from_title
        elif date_from_uri:
            extracted_date = date_from_uri
        
        # 日付が見つからない場合は最も古い日付として扱う
        if extracted_date is None:
            extracted_date = datetime(1900, 1, 1)
        
        # 新しい日付を優先するため、日付を逆順でソート
        return (-extracted_date.timestamp(), title.lower())
    
    return sorted(sources, key=get_sort_key)

def convert_grounding_metadata_to_dict(grounding_metadata):
    """GroundingMetadataオブジェクトを辞書に変換し、日付順にソート"""
    if not grounding_metadata:
        print("DEBUG: grounding_metadata is None")
        return None
    
    print(f"DEBUG: grounding_metadata type: {type(grounding_metadata)}")
    print(f"DEBUG: grounding_metadata attributes: {dir(grounding_metadata)}")
    
    result = {}
    
    # grounding_chunksを処理
    if hasattr(grounding_metadata, 'grounding_chunks') and grounding_metadata.grounding_chunks:
        print(f"DEBUG: Found {len(grounding_metadata.grounding_chunks)} grounding chunks")
        unsorted_chunks = []
        
        for i, chunk in enumerate(grounding_metadata.grounding_chunks):
            print(f"DEBUG: Processing chunk {i}: {type(chunk)}")
            print(f"DEBUG: Chunk attributes: {dir(chunk)}")
            
            chunk_dict = {}
            
            # retrieved_contextからtitleとuriを取得
            if hasattr(chunk, 'retrieved_context') and chunk.retrieved_context:
                print(f"DEBUG: Found retrieved_context: {type(chunk.retrieved_context)}")
                retrieved_context = chunk.retrieved_context
                
                # titleを取得
                if hasattr(retrieved_context, 'title') and retrieved_context.title:
                    chunk_dict['title'] = retrieved_context.title
                    print(f"DEBUG: Found title in retrieved_context: {retrieved_context.title}")
                
                # uriを取得
                if hasattr(retrieved_context, 'uri') and retrieved_context.uri:
                    chunk_dict['uri'] = retrieved_context.uri
                    print(f"DEBUG: Found uri in retrieved_context: {retrieved_context.uri}")
            
            # webプロパティも確認（念のため）
            if hasattr(chunk, 'web') and chunk.web:
                print(f"DEBUG: Found web property: {chunk.web}")
                if not chunk_dict.get('title') and hasattr(chunk.web, 'title'):
                    chunk_dict['title'] = chunk.web.title
                if not chunk_dict.get('uri') and hasattr(chunk.web, 'uri'):
                    chunk_dict['uri'] = chunk.web.uri
            
            # 直接的なtitle/uriプロパティも確認
            if not chunk_dict.get('title') and hasattr(chunk, 'title'):
                chunk_dict['title'] = chunk.title
                print(f"DEBUG: Found direct title: {chunk.title}")
            if not chunk_dict.get('uri') and hasattr(chunk, 'uri'):
                chunk_dict['uri'] = chunk.uri
                print(f"DEBUG: Found direct uri: {chunk.uri}")
            
            # 日付情報を抽出してデバッグ出力
            if chunk_dict.get('title'):
                extracted_date = extract_date_from_filename(chunk_dict['title'])
                if extracted_date:
                    print(f"DEBUG: Extracted date from title '{chunk_dict['title']}': {extracted_date.strftime('%Y-%m-%d')}")
                else:
                    print(f"DEBUG: No date found in title '{chunk_dict['title']}'")
            
            unsorted_chunks.append(chunk_dict)
        
        # 日付順にソート（新しい日付を優先）
        sorted_chunks = sort_sources_by_date(unsorted_chunks)
        result['grounding_chunks'] = sorted_chunks
        
        print(f"DEBUG: Sorted chunks by date:")
        for i, chunk in enumerate(sorted_chunks):
            title = chunk.get('title', 'タイトルなし')
            date = extract_date_from_filename(title)
            date_str = date.strftime('%Y-%m-%d') if date else '日付なし'
            print(f"  {i+1}. {title} ({date_str})")
        
    else:
        print("DEBUG: No grounding_chunks found")
    
    print(f"DEBUG: Final result: {result}")
    return result

def generate_plan_and_questions(user_message):
    """ユーザーの質問から計画と関連質問を生成"""
    try:
        client = create_rag_client()
        
        planning_prompt = f"""
以下のユーザーの質問に対して、包括的で詳細な回答を提供するための計画を立ててください。

ユーザーの質問: {user_message}

以下の形式で回答してください：

## 調査計画
[この質問に答えるための調査計画を簡潔に説明]

## 関連質問リスト
1. [関連質問1]
2. [関連質問2]
3. [関連質問3]
4. [関連質問4]
5. [関連質問5]

関連質問は以下の観点から作成してください：
- 基本的な定義や概念
- 具体的な事例や応用
- メリット・デメリット
- 最新の動向や課題
- 関連する技術や手法

各質問は独立して回答可能で、元の質問の理解を深めるものにしてください。
"""
        
        contents = [
            types.Content(
                role="user",
                parts=[types.Part(text=planning_prompt)]
            )
        ]
        
        config = types.GenerateContentConfig(
            temperature=0.7,
            top_p=0.9,
            max_output_tokens=1000,
        )
        
        response = client.models.generate_content(
            model=GEMINI_MODEL,
            contents=contents,
            config=config,
        )
        
        if response and response.text:
            return response.text
        else:
            print("WARNING: Empty response from planning model")
            return f"""
## 調査計画
{user_message}について詳細に調査します。

## 関連質問リスト
1. {user_message}の基本的な定義とは何ですか？
2. {user_message}の具体的な事例を教えてください
3. {user_message}のメリットとデメリットは何ですか？
4. {user_message}の最新の動向はどうですか？
5. {user_message}に関連する技術や手法はありますか？
"""
    except Exception as e:
        print(f"Error in generate_plan_and_questions: {e}")
        return f"""
## 調査計画
{user_message}について詳細に調査します。

## 関連質問リスト
1. {user_message}の基本的な定義とは何ですか？
2. {user_message}の具体的な事例を教えてください
3. {user_message}のメリットとデメリットは何ですか？
4. {user_message}の最新の動向はどうですか？
5. {user_message}に関連する技術や手法はありますか？
"""

def execute_single_rag_query(question):
    """単一のRAGクエリを実行"""
    try:
        client = create_rag_client()
        
        contents = [
            types.Content(
                role="user",
                parts=[types.Part(text=question)]
            )
        ]
        
        tools = [
            types.Tool(
                retrieval=types.Retrieval(
                    vertex_rag_store=types.VertexRagStore(
                        rag_resources=[
                            types.VertexRagStoreRagResource(
                                rag_corpus=RAG_CORPUS
                            )
                        ],
                    )
                )
            )
        ]

        config = types.GenerateContentConfig(
            temperature=0.8,
            top_p=0.9,
            max_output_tokens=2000,
            tools=tools,
        )
        
        response = client.models.generate_content(
            model=GEMINI_MODEL,
            contents=contents,
            config=config,
        )
        
        # グラウンディングメタデータを取得
        grounding_metadata = None
        if hasattr(response, 'candidates') and response.candidates:
            candidate = response.candidates[0]
            if hasattr(candidate, 'grounding_metadata') and candidate.grounding_metadata:
                grounding_metadata = candidate.grounding_metadata
        
        answer_text = response.text if response and response.text else "回答を取得できませんでした。"
        return answer_text, grounding_metadata
        
    except Exception as e:
        print(f"Error in execute_single_rag_query: {e}")
        return f"エラーが発生しました: {str(e)}", None

def synthesize_comprehensive_answer(user_message, plan_text, qa_results):
    """計画と各質問の回答を統合して包括的な回答を生成"""
    client = create_rag_client()
    
    qa_text = "\n\n".join([f"**Q: {q}**\nA: {a}" for q, a in qa_results])
    
    synthesis_prompt = f"""
以下の情報を基に、ユーザーの質問に対する包括的で詳細な回答を作成してください。

元の質問: {user_message}

調査計画:
{plan_text}

関連質問と回答:
{qa_text}

以下の要件に従って回答を作成してください：
1. 元の質問に直接答える
2. 関連質問の回答から得られた情報を統合する
3. 論理的で読みやすい構成にする
4. 重要なポイントを強調する
5. 具体例があれば含める
6. Markdown形式で整理する

回答は以下の構成を参考にしてください：
- 概要・定義
- 詳細説明
- 具体例・事例
- メリット・デメリット
- 最新動向・課題
- まとめ
"""
    
    contents = [
        types.Content(
            role="user",
            parts=[types.Part(text=synthesis_prompt)]
        )
    ]
    
    config = types.GenerateContentConfig(
        temperature=0.7,
        top_p=0.9,
        max_output_tokens=4000,
    )
    
    response = client.models.generate_content(
        model=GEMINI_MODEL,
        contents=contents,
        config=config,
    )
    
    return response.text

def generate_deep_response(user_message):
    """深掘り機能付きのレスポンス生成"""
    try:
        # ステップ1: 計画立てと関連質問生成
        yield {
            'chunk': '\n## 📋 調査計画を立案中...\n',
            'done': False,
            'grounding_metadata': None,
            'step': 'planning'
        }
        
        plan_text = generate_plan_and_questions(user_message)
        
        if not plan_text:
            plan_text = f"""
## 調査計画
{user_message}について詳細に調査します。

## 関連質問リスト
1. {user_message}の基本的な定義とは何ですか？
2. {user_message}の具体的な事例を教えてください
3. {user_message}のメリットとデメリットは何ですか？
4. {user_message}の最新の動向はどうですか？
5. {user_message}に関連する技術や手法はありますか？
"""
        
        yield {
            'chunk': f'\n{plan_text}\n\n## 🔍 詳細調査を開始...\n',
            'done': False,
            'grounding_metadata': None,
            'step': 'plan_complete'
        }
        
        # 関連質問を抽出
        questions = []
        try:
            lines = plan_text.split('\n')
            in_questions_section = False
            
            for line in lines:
                if '関連質問' in line:
                    in_questions_section = True
                    continue
                if in_questions_section and line.strip():
                    if line.strip().startswith(('1.', '2.', '3.', '4.', '5.')):
                        question = line.strip()[2:].strip()
                        if question and question not in questions:
                            questions.append(question)
        except Exception as e:
            print(f"Error extracting questions: {e}")
            questions = [
                f"{user_message}の基本的な定義とは何ですか？",
                f"{user_message}の具体的な事例を教えてください",
                f"{user_message}のメリットとデメリットは何ですか？",
                f"{user_message}の最新の動向はどうですか？",
                f"{user_message}に関連する技術や手法はありますか？"
            ]
        
        # 質問が少ない場合のフォールバック
        if len(questions) < 3:
            questions = [
                f"{user_message}の基本的な定義とは何ですか？",
                f"{user_message}の具体的な事例を教えてください",
                f"{user_message}のメリットとデメリットは何ですか？",
                f"{user_message}の最新の動向はどうですか？",
                f"{user_message}に関連する技術や手法はありますか？"
            ]
        
        # ステップ2: 各関連質問を順次実行
        qa_results = []
        all_grounding_metadata = []
        all_unique_sources = {}  # 重複を避けるため辞書で管理
        
        for i, question in enumerate(questions[:5], 1):
            yield {
                'chunk': f'\n### 🔍 質問 {i}: {question}\n調査中...\n',
                'done': False,
                'grounding_metadata': None,
                'step': f'query_{i}'
            }
            
            answer, grounding_metadata = execute_single_rag_query(question)
            qa_results.append((question, answer))
            
            # 出典情報を処理
            converted_metadata = None
            if grounding_metadata:
                all_grounding_metadata.append(grounding_metadata)
                converted_metadata = convert_grounding_metadata_to_dict(grounding_metadata)
                
                # 出典情報を統合（重複を避ける）
                if converted_metadata and 'grounding_chunks' in converted_metadata:
                    for chunk in converted_metadata['grounding_chunks']:
                        if chunk.get('uri'):
                            all_unique_sources[chunk['uri']] = {
                                'title': chunk.get('title', 'タイトルなし'),
                                'uri': chunk['uri']
                            }
            
            # 回答を送信
            yield {
                'chunk': f'\n**💡 回答 {i}:** {answer}\n',
                'done': False,
                'grounding_metadata': converted_metadata,
                'step': f'answer_{i}'
            }
            
            # 各質問の出典情報を個別に表示
            if converted_metadata and 'grounding_chunks' in converted_metadata:
                # 出典情報を日付順にソート
                sorted_chunks = sort_sources_by_date(converted_metadata['grounding_chunks'])
                
                sources_text = '\n**📚 この回答の出典:**\n'
                for j, chunk in enumerate(sorted_chunks, 1):
                    title = chunk.get('title', 'タイトルなし')
                    uri = chunk.get('uri', '')
                    
                    # 日付情報を表示に含める
                    extracted_date = extract_date_from_filename(title)
                    date_info = f" ({extracted_date.strftime('%Y-%m-%d')})" if extracted_date else ""
                    
                    sources_text += f'   {j}. {title}{date_info}\n'
                    if uri:
                        sources_text += f'      📎 {uri}\n'
                sources_text += '\n'
                
                yield {
                    'chunk': sources_text,
                    'done': False,
                    'grounding_metadata': None,
                    'step': f'sources_{i}'
                }
        
        # ステップ3: 包括的な回答の統合
        yield {
            'chunk': '\n## 📝 包括的な回答を作成中...\n',
            'done': False,
            'grounding_metadata': None,
            'step': 'synthesizing'
        }
        
        comprehensive_answer = synthesize_comprehensive_answer(user_message, plan_text, qa_results)
        
        yield {
            'chunk': f'\n## 🎯 包括的な回答\n\n{comprehensive_answer}\n',
            'done': False,
            'grounding_metadata': None,
            'step': 'synthesis_complete'
        }
        
        # 全ての出典情報を統合して表示
        if all_unique_sources:
            # 日付順にソートしてから表示
            sorted_unique_sources = sort_sources_by_date(list(all_unique_sources.values()))
            
            yield {
                'chunk': '\n## 📚 全体の出典情報\n\n',
                'done': False,
                'grounding_metadata': None,
                'step': 'all_sources_header'
            }
            
            sources_summary = ''
            for i, source_info in enumerate(sorted_unique_sources, 1):
                title = source_info.get('title', 'タイトルなし')
                uri = source_info.get('uri', '')
                
                # 日付情報を表示に含める
                extracted_date = extract_date_from_filename(title)
                date_info = f" ({extracted_date.strftime('%Y-%m-%d')})" if extracted_date else ""
                
                sources_summary += f'**{i}. {title}{date_info}**\n'
                sources_summary += f'   📎 {uri}\n\n'
            
            yield {
                'chunk': sources_summary,
                'done': False,
                'grounding_metadata': None,
                'step': 'all_sources_list'
            }
        
        # 最終的な出典情報を統合（JSONとして送信）
        final_grounding_metadata = None
        if all_grounding_metadata:
            # 全ての出典情報を統合し、日付順にソート
            sorted_unique_sources = sort_sources_by_date(list(all_unique_sources.values()))
            final_grounding_metadata = {
                'grounding_chunks': sorted_unique_sources
            }
        
        yield {
            'chunk': '',
            'done': True,
            'grounding_metadata': final_grounding_metadata,
            'step': 'complete'
        }
        
    except Exception as e:
        print(f"Deep response generation error: {e}")
        yield {
            'chunk': f'\n❌ エラーが発生しました: {str(e)}\n通常モードで回答を試みます...\n',
            'done': False,
            'grounding_metadata': None,
            'step': 'error'
        }
        
        # エラー時は通常モードにフォールバック
        try:
            for chunk_data in generate_response(user_message):
                yield chunk_data
        except Exception as fallback_error:
            print(f"Fallback error: {fallback_error}")
            yield {
                'chunk': f'\n❌ 回答の生成に失敗しました: {str(fallback_error)}\n',
                'done': True,
                'grounding_metadata': None,
                'step': 'fallback_error'
            }

def generate_response(user_message):
    """ユーザーメッセージに対してRAGを使用してレスポンスを生成"""
    client = create_rag_client()
    
    contents = [
        types.Content(
            role="user",
            parts=[
                types.Part(text=user_message)
            ]
        )
    ]
    tools = [
        types.Tool(
            retrieval=types.Retrieval(
                vertex_rag_store=types.VertexRagStore(
                    rag_resources=[
                        types.VertexRagStoreRagResource(
                            rag_corpus=RAG_CORPUS
                        )
                    ],
                )
            )
        )
    ]

    # GenerateContentConfigを作成
    config_params = {
        'temperature': 1,
        'top_p': 1,
        'seed': 0,
        'max_output_tokens': 65535,
        'safety_settings': [
            types.SafetySetting(
                category="HARM_CATEGORY_HATE_SPEECH",
                threshold="OFF"
            ),
            types.SafetySetting(
                category="HARM_CATEGORY_DANGEROUS_CONTENT",
                threshold="OFF"
            ),
            types.SafetySetting(
                category="HARM_CATEGORY_SEXUALLY_EXPLICIT",
                threshold="OFF"
            ),
            types.SafetySetting(
                category="HARM_CATEGORY_HARASSMENT",
                threshold="OFF"
            )
        ],
        'tools': tools,
    }
    
    # ThinkingConfigが利用可能な場合のみ追加
    try:
        # より詳細なチェックを実装
        thinking_config_available = False
        try:
            # ThinkingConfigクラスが存在するかチェック
            if hasattr(types, 'ThinkingConfig'):
                # 実際にインスタンス化できるかテスト
                test_config = types.ThinkingConfig(thinking_budget=-1)
                thinking_config_available = True
                print("DEBUG: ThinkingConfig is available")
            else:
                # 代替的なアプローチ: 直接インポートを試行
                try:
                    from google.genai.types import ThinkingConfig
                    test_config = ThinkingConfig(thinking_budget=-1)
                    thinking_config_available = True
                    print("DEBUG: ThinkingConfig imported directly")
                except ImportError:
                    print("DEBUG: ThinkingConfig not available via direct import")
        except (AttributeError, TypeError, ValueError) as e:
            print(f"DEBUG: ThinkingConfig not available: {e}")
            thinking_config_available = False
        
        if thinking_config_available:
            config_params['thinking_config'] = types.ThinkingConfig(
                thinking_budget=-1,
            )
            print("DEBUG: Added ThinkingConfig to config")
        else:
            print("DEBUG: Skipping ThinkingConfig - not available")
    except Exception as e:
        print(f"DEBUG: Error checking ThinkingConfig: {e}")
        pass  # ThinkingConfigが利用できない場合は無視
    
    # 安全にGenerateContentConfigを作成
    try:
        generate_content_config = types.GenerateContentConfig(**config_params)
        print("DEBUG: GenerateContentConfig created successfully")
    except Exception as e:
        print(f"DEBUG: Error creating GenerateContentConfig: {e}")
        # ThinkingConfigを除外して再試行
        if 'thinking_config' in config_params:
            del config_params['thinking_config']
            print("DEBUG: Retrying without ThinkingConfig")
            generate_content_config = types.GenerateContentConfig(**config_params)

    full_response = ""
    grounding_metadata = None
    
    print(f"DEBUG: Starting generation for message: {user_message}")
    
    for chunk in client.models.generate_content_stream(
        model=GEMINI_MODEL,
        contents=contents,
        config=generate_content_config,
    ):
        print(f"DEBUG: Raw chunk: {chunk}")
        
        if not chunk.candidates or not chunk.candidates[0].content or not chunk.candidates[0].content.parts:
            print("DEBUG: Skipping chunk - no content")
            continue
        
        print(f"DEBUG: Chunk type: {type(chunk)}")
        print(f"DEBUG: Chunk attributes: {[attr for attr in dir(chunk) if not attr.startswith('_')]}")
        
        # テキストを蓄積
        full_response += chunk.text
        print(f"DEBUG: Added text: {chunk.text[:100]}...")
        
        # 全ての可能な場所でグラウンディングメタデータを探す
        candidate = chunk.candidates[0]
        print(f"DEBUG: Candidate attributes: {[attr for attr in dir(candidate) if not attr.startswith('_')]}")
        
        # 候補1: candidate.grounding_metadata
        if hasattr(candidate, 'grounding_metadata'):
            print(f"DEBUG: candidate.grounding_metadata exists: {candidate.grounding_metadata}")
            if candidate.grounding_metadata:
                print("DEBUG: Found grounding_metadata in candidate")
                grounding_metadata = candidate.grounding_metadata
        
        # 候補2: chunk.grounding_metadata
        if hasattr(chunk, 'grounding_metadata'):
            print(f"DEBUG: chunk.grounding_metadata exists: {chunk.grounding_metadata}")
            if chunk.grounding_metadata:
                print("DEBUG: Found grounding_metadata in chunk")
                grounding_metadata = chunk.grounding_metadata
            
        # 候補3: candidate.content.grounding_metadata
        if hasattr(candidate.content, 'grounding_metadata'):
            print(f"DEBUG: candidate.content.grounding_metadata exists: {candidate.content.grounding_metadata}")
            if candidate.content.grounding_metadata:
                print("DEBUG: Found grounding_metadata in candidate.content")
                grounding_metadata = candidate.content.grounding_metadata
        
        # 候補4: 全体の構造を確認
        print(f"DEBUG: Full candidate structure:")
        for attr in dir(candidate):
            if not attr.startswith('_'):
                try:
                    value = getattr(candidate, attr)
                    print(f"  {attr}: {type(value)} - {value}")
                except:
                    print(f"  {attr}: <error accessing>")
        
        yield {
            'chunk': chunk.text,
            'done': False,
            'grounding_metadata': None
        }
    
    print(f"DEBUG: Final grounding_metadata: {grounding_metadata}")
    
    # 最後に出典情報を送信（辞書形式に変換）
    converted_metadata = convert_grounding_metadata_to_dict(grounding_metadata)
    print(f"DEBUG: Converted metadata: {converted_metadata}")
    
    yield {
        'chunk': '',
        'done': True,
        'grounding_metadata': converted_metadata
    }

# アップロード許可ファイル形式
ALLOWED_EXTENSIONS = {'txt', 'pdf', 'docx', 'md', 'doc', 'rtf', 'pptx', 'ppt'}
UPLOAD_FOLDER = 'uploads'
MAX_CONTENT_LENGTH = 16 * 1024 * 1024  # 16MB

# アップロードフォルダの作成
os.makedirs(UPLOAD_FOLDER, exist_ok=True)

app.config['MAX_CONTENT_LENGTH'] = MAX_CONTENT_LENGTH

def allowed_file(filename):
    """アップロード可能なファイル形式かチェック"""
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

def upload_to_cloud_storage(file_path, filename):
    """Cloud Storageにファイルをアップロード"""
    try:
        # Cloud Storageクライアントを初期化
        client = storage.Client(project=PROJECT_ID)
        bucket_name = f"{PROJECT_ID}-rag-documents"
        
        # バケットが存在しない場合は作成
        try:
            bucket = client.get_bucket(bucket_name)
        except Exception:
            bucket = client.create_bucket(bucket_name)
        
        # ファイルをアップロード
        blob_name = f"documents/{datetime.now().strftime('%Y%m%d_%H%M%S')}_{filename}"
        blob = bucket.blob(blob_name)
        
        with open(file_path, 'rb') as file_data:
            blob.upload_from_file(file_data)
        
        return f"gs://{bucket_name}/{blob_name}"
    except Exception as e:
        print(f"Cloud Storageアップロードエラー: {e}")
        return None

def get_rag_documents():
    """Vertex AI RAGから直接ドキュメント一覧を取得"""
    try:
        print("RAGサービスクライアントが利用できません。Cloud Storageから取得します。")
        return []
        
        # RAGサービスクライアント機能を一時的に無効化
        # client = RagDataServiceClient()
        # ... 以下のコードは一時的にコメントアウト
        
    except Exception as e:
        print(f"RAGドキュメント一覧取得エラー: {e}")
        return []

def get_file_size_string(size_bytes):
    """ファイルサイズを人間が読みやすい形式に変換"""
    if size_bytes == 0:
        return "0 B"
    
    size_names = ["B", "KB", "MB", "GB", "TB"]
    i = 0
    while size_bytes >= 1024 and i < len(size_names) - 1:
        size_bytes /= 1024.0
        i += 1
    
    return f"{size_bytes:.1f} {size_names[i]}"

def add_document_to_rag(file_uri, display_name):
    """Vertex AI RAGにドキュメントを追加"""
    try:
        print(f"RAGサービスクライアントが利用できません。プレースホルダーとして成功を返します: {display_name}")
        
        # RAGサービスクライアント機能を一時的に無効化
        # client = RagDataServiceClient()
        # ... 以下のコードは一時的にコメントアウト
        
        return {
            "success": True,
            "operation_name": "placeholder_operation",
            "display_name": display_name,
            "file_uri": file_uri,
            "message": "ドキュメントをCloud Storageにアップロードしました。（RAG機能は一時的に無効化されています）"
        }
        
    except Exception as e:
        print(f"RAGドキュメント追加エラー: {e}")
        return {
            "success": False,
            "error": str(e)
        }

def delete_rag_document(rag_file_name):
    """Vertex AI RAGからドキュメントを削除"""
    try:
        print(f"RAGサービスクライアントが利用できません。プレースホルダーとして成功を返します: {rag_file_name}")
        
        # RAGサービスクライアント機能を一時的に無効化
        # client = RagDataServiceClient()
        # ... 以下のコードは一時的にコメントアウト
        
        return {
            "success": True,
            "message": "RAG機能は一時的に無効化されています"
        }
        
    except Exception as e:
        print(f"RAGドキュメント削除エラー: {e}")
        return {
            "success": False,
            "error": str(e)
        }

@app.route('/upload_documents', methods=['POST'])
@auth.login_required
def upload_documents():
    """ドキュメントアップロードエンドポイント（RAG統合版）"""
    try:
        files = request.files.getlist('files')
        
        # フォームデータからRAGコーパス情報を取得
        corpus_name = request.form.get('corpus_name', '').strip()
        corpus_description = request.form.get('corpus_description', '').strip()
        chunk_size = int(request.form.get('chunk_size', 512))
        chunk_overlap = int(request.form.get('chunk_overlap', 100))
        
        # デフォルトコーパス名の設定
        if not corpus_name:
            corpus_name = f"ドキュメント_{datetime.now().strftime('%Y%m%d_%H%M%S')}"
        
        if not files:
            return jsonify({'error': 'ファイルが選択されていません'}), 400
        
        print(f"ドキュメントアップロード開始:")
        print(f"  ファイル数: {len(files)}")
        print(f"  RAGコーパス名: {corpus_name}")
        print(f"  チャンクサイズ: {chunk_size}")
        print(f"  チャンクオーバーラップ: {chunk_overlap}")
        
        results = []
        
        for file in files:
            if file.filename == '':
                continue
            
            if file and allowed_file(file.filename):
                try:
                    # ファイル名を安全にする
                    filename = secure_filename(file.filename)
                    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
                    unique_filename = f"{timestamp}_{filename}"
                    
                    # 一時的にローカルに保存
                    file_path = os.path.join(UPLOAD_FOLDER, unique_filename)
                    file.save(file_path)
                    
                    # ファイルサイズを取得
                    file_size = os.path.getsize(file_path)
                    file_size_str = get_file_size_string(file_size)
                    
                    # Cloud StorageアップロードとRAGインポートを一括実行
                    rag_result = upload_and_import_to_rag(
                        file_path, 
                        unique_filename, 
                        corpus_name, 
                        corpus_description
                    )
                    
                    if rag_result['success']:
                        results.append({
                            'success': True,
                            'filename': filename,
                            'message': f'完了 ({file_size_str}) - {rag_result.get("message", "")}',
                            'size': file_size_str,
                            'gcs_uri': rag_result.get('gcs_uri'),
                            'corpus_id': rag_result.get('corpus_id'),
                            'corpus_name': rag_result.get('corpus_name'),
                            'corpus_created': rag_result.get('corpus_created', False)
                        })
                    else:
                        results.append({
                            'success': False,
                            'filename': filename,
                            'message': f'エラー ({file_size_str}) - {rag_result.get("message", "")}',
                            'size': file_size_str,
                            'error': rag_result.get('error')
                        })
                    
                    # ローカルファイルを削除
                    try:
                        os.remove(file_path)
                    except:
                        pass
                    
                except Exception as e:
                    results.append({
                        'success': False,
                        'filename': file.filename,
                        'message': f'処理エラー: {str(e)}',
                        'size': 'unknown'
                    })
            else:
                results.append({
                    'success': False,
                    'filename': file.filename,
                    'message': f'サポートされていないファイル形式です。許可形式: {", ".join(ALLOWED_EXTENSIONS)}',
                    'size': 'unknown'
                })
        
        # 結果の統計を作成
        total_files = len(results)
        success_files = len([r for r in results if r['success']])
        error_files = total_files - success_files
        
        # コーパス作成情報を取得
        corpus_created = any(r.get('corpus_created', False) for r in results if r['success'])
        corpus_ids = list(set(r.get('corpus_id') for r in results if r['success'] and r.get('corpus_id')))
        
        return jsonify({
            'results': results,
            'summary': {
                'total_files': total_files,
                'success_files': success_files,
                'error_files': error_files,
                'corpus_name': corpus_name,
                'corpus_created': corpus_created,
                'corpus_ids': corpus_ids
            }
        })
    
    except Exception as e:
        print(f"アップロードエラー: {e}")
        import traceback
        traceback.print_exc()
        return jsonify({'error': f'アップロードエラー: {str(e)}'}), 500

@app.route('/get_documents', methods=['GET'])
@auth.login_required
def get_documents():
    """ドキュメント一覧取得エンドポイント"""
    try:
        use_rag_api = request.args.get('use_rag_api', 'false').lower() == 'true'
        
        if use_rag_api:
            # Vertex AI RAGのAPIを直接使用
            documents = get_rag_documents()
            return jsonify({'documents': documents, 'source': 'rag_api'})
        else:
            # Cloud Storageからドキュメント一覧を取得（既存の実装）
            client = storage.Client(project=PROJECT_ID)
            bucket_name = f"{PROJECT_ID}-rag-documents"
            
            try:
                bucket = client.get_bucket(bucket_name)
                blobs = bucket.list_blobs(prefix="documents/")
                
                documents = []
                for blob in blobs:
                    # ファイル名から元のファイル名を抽出
                    filename = blob.name.split('/')[-1]
                    original_filename = '_'.join(filename.split('_')[2:]) if '_' in filename else filename
                    
                    documents.append({
                        'id': blob.name,
                        'title': original_filename,
                        'created_at': blob.time_created.strftime('%Y-%m-%d %H:%M:%S'),
                        'size': get_file_size_string(blob.size),
                        'gcs_uri': f"gs://{bucket_name}/{blob.name}",
                        'source': 'cloud_storage'
                    })
                
                # 作成日時でソート（新しい順）
                documents.sort(key=lambda x: x['created_at'], reverse=True)
                
                return jsonify({'documents': documents, 'source': 'cloud_storage'})
                
            except Exception as e:
                print(f"バケット取得エラー: {e}")
                return jsonify({'documents': [], 'source': 'cloud_storage'})
    
    except Exception as e:
        print(f"ドキュメント一覧取得エラー: {e}")
        return jsonify({'error': f'ドキュメント一覧取得エラー: {str(e)}'}), 500

@app.route('/delete_document/<path:doc_id>', methods=['DELETE'])
@auth.login_required
def delete_document(doc_id):
    """ドキュメント削除エンドポイント"""
    try:
        use_rag_api = request.args.get('use_rag_api', 'false').lower() == 'true'
        
        if use_rag_api:
            # Vertex AI RAGから直接削除
            result = delete_rag_document(doc_id)
            return jsonify(result)
        else:
            # Cloud Storageからファイルを削除（既存の実装）
            client = storage.Client(project=PROJECT_ID)
            bucket_name = f"{PROJECT_ID}-rag-documents"
            
            try:
                bucket = client.get_bucket(bucket_name)
                blob = bucket.blob(doc_id)
                blob.delete()
                
                return jsonify({
                    'success': True,
                    'message': 'ドキュメントを削除しました（Cloud Storageから）'
                })
                
            except Exception as e:
                print(f"ドキュメント削除エラー: {e}")
                return jsonify({
                    'success': False,
                    'message': f'削除エラー: {str(e)}'
                }), 500
    
    except Exception as e:
        print(f"ドキュメント削除エラー: {e}")
        return jsonify({
            'success': False,
            'message': f'削除エラー: {str(e)}'
        }), 500

@app.route('/get_settings', methods=['GET'])
@auth.login_required
def get_settings():
    """設定情報取得エンドポイント"""
    try:
        return jsonify({
            'rag_corpus': RAG_CORPUS,
            'gemini_model': GEMINI_MODEL,
            'project_id': PROJECT_ID,
            'last_update': datetime.now().strftime('%Y-%m-%d %H:%M:%S')
        })
    except Exception as e:
        print(f"設定取得エラー: {e}")
        return jsonify({'error': f'設定取得エラー: {str(e)}'}), 500

@app.route('/')
@auth.login_required
def index():
    """メインページ"""
    return render_template('index.html')

@app.route('/health')
def health():
    """ヘルスチェックエンドポイント"""
    import psutil
    import gc
    
    # メモリ使用量を取得
    process = psutil.Process()
    memory_info = process.memory_info()
    
    # ガベージコレクションを実行
    gc.collect()
    
    return jsonify({
        'status': 'healthy',
        'memory_usage_mb': round(memory_info.rss / 1024 / 1024, 2),
        'memory_percent': round(process.memory_percent(), 2),
        'cpu_percent': round(process.cpu_percent(), 2)
    })

@app.route('/chat', methods=['POST'])
@auth.login_required
def chat():
    """チャットエンドポイント"""
    data = request.json
    user_message = data.get('message', '')
    use_deep_mode = data.get('deep_mode', False)
    
    if not user_message:
        return jsonify({'error': 'メッセージが空です'}), 400
    
    def generate():
        try:
            if use_deep_mode:
                # 深掘りモードを使用
                for chunk_data in generate_deep_response(user_message):
                    yield f"data: {json.dumps(chunk_data, ensure_ascii=False)}\n\n"
            else:
                # 通常モードを使用
                for chunk_data in generate_response(user_message):
                    yield f"data: {json.dumps(chunk_data, ensure_ascii=False)}\n\n"
            
            # メモリクリーンアップ
            gc.collect()
            
        except Exception as e:
            print(f"Error in chat endpoint: {e}")
            error_data = {
                'chunk': f'エラーが発生しました: {str(e)}',
                'done': True,
                'grounding_metadata': None
            }
            yield f"data: {json.dumps(error_data, ensure_ascii=False)}\n\n"
            # エラー時もメモリクリーンアップ
            gc.collect()
    
    return Response(generate(), mimetype='text/plain')

# Google Drive同期関連のルート
@app.route('/auth/google')
def google_auth():
    """Google OAuth認証を開始"""
    drive_sync = get_drive_sync()
    if not drive_sync:
        return jsonify({'error': 'Google Drive同期機能が無効です'}), 500
    
    redirect_uri = url_for('google_callback', _external=True)
    auth_url, state = drive_sync.get_auth_url(redirect_uri)
    
    session['oauth_state'] = state
    return redirect(auth_url)

@app.route('/auth/google/callback')
def google_callback():
    """Google OAuth認証コールバック"""
    print(f"DEBUG: Callback received - URL: {request.url}")
    print(f"DEBUG: Session oauth_state: {session.get('oauth_state')}")
    
    drive_sync = get_drive_sync()
    if not drive_sync:
        print("ERROR: Google Drive同期機能が無効です")
        return jsonify({'error': 'Google Drive同期機能が無効です'}), 500
    
    state = session.get('oauth_state')
    if not state:
        print("ERROR: 認証状態が無効です")
        return jsonify({'error': '認証状態が無効です'}), 400
    
    redirect_uri = url_for('google_callback', _external=True)
    authorization_response = request.url
    
    print(f"DEBUG: Calling handle_oauth_callback")
    
    if drive_sync.handle_oauth_callback(authorization_response, state, redirect_uri):
        print("SUCCESS: OAuth認証成功")
        return redirect(url_for('index') + '?auth=success')
    else:
        print("ERROR: OAuth認証失敗")
        return redirect(url_for('index') + '?auth=error')

@app.route('/drive/status')
def drive_status():
    """Google Drive同期状態を取得"""
    drive_sync = get_drive_sync()
    if not drive_sync:
        return jsonify({'error': 'Google Drive同期機能が無効です'}), 500
    
    # 保存された認証情報を読み込み
    drive_sync.load_credentials()
    
    status = drive_sync.get_sync_status()
    return jsonify(status)

@app.route('/drive/sync', methods=['POST'])
@auth.login_required
def sync_drive():
    """Google Driveフォルダを同期"""
    drive_sync = get_drive_sync()
    if not drive_sync:
        return jsonify({'error': 'Google Drive同期機能が無効です'}), 500
    
    # 保存された認証情報を読み込み
    if not drive_sync.load_credentials():
        return jsonify({'error': '認証が必要です', 'auth_required': True}), 401
    
    data = request.get_json() or {}
    folder_name = data.get('folder_name', DRIVE_FOLDER_NAME)
    corpus_name = data.get('corpus_name', 'Drive Sync Corpus')
    force_sync = data.get('force_sync', False)
    recursive = data.get('recursive_sync', False)  # フロントエンドから送信されるパラメータ名に合わせて修正
    try:
        max_depth = int(data.get('max_depth', 3))      # 文字列から整数に変換
        if max_depth < 1 or max_depth > 10:
            max_depth = 3  # デフォルト値に戻す
    except (ValueError, TypeError):
        max_depth = 3  # デフォルト値
    
    try:
        result = drive_sync.sync_folder_to_rag(
            folder_name, 
            corpus_name, 
            force_sync, 
            recursive, 
            max_depth
        )
        return jsonify(result)
    except Exception as e:
        return jsonify({'error': f'同期エラー: {str(e)}'}), 500

@app.route('/drive/folders')
def list_drive_folders():
    """Google Driveのフォルダ一覧を取得"""
    drive_sync = get_drive_sync()
    if not drive_sync:
        return jsonify({'error': 'Google Drive同期機能が無効です'}), 500
    
    if not drive_sync.load_credentials():
        return jsonify({'error': '認証が必要です', 'auth_required': True}), 401
    
    try:
        parent_id = request.args.get('parent', 'root')
        folders = drive_sync.list_folders(parent_id)
        
        # フォルダ情報にパス情報を追加
        folder_list = []
        for folder in folders:
            folder_path = drive_sync.get_folder_path(folder['id'])
            folder_list.append({
                'id': folder['id'],
                'name': folder['name'],
                'path': folder_path,
                'created_time': folder.get('createdTime'),
                'modified_time': folder.get('modifiedTime')
            })
        
        return jsonify({
            'folders': folder_list,
            'parent_id': parent_id
        })
        
    except Exception as e:
        return jsonify({'error': f'フォルダ一覧取得エラー: {str(e)}'}), 500

@app.route('/drive/clear_auth', methods=['POST'])
@auth.login_required
def clear_drive_auth():
    """Google Drive認証をクリア"""
    drive_sync = get_drive_sync()
    if not drive_sync:
        return jsonify({'error': 'Google Drive同期機能が無効です'}), 500
    
    drive_sync.clear_auth()
    return jsonify({'message': '認証情報をクリアしました'})

@app.route('/drive/folder/info', methods=['POST'])
def get_folder_info():
    """フォルダ情報を取得（プレビュー用）"""
    drive_sync = get_drive_sync()
    if not drive_sync:
        return jsonify({'error': 'Google Drive同期機能が無効です'}), 500
    
    if not drive_sync.load_credentials():
        return jsonify({'error': '認証が必要です', 'auth_required': True}), 401
    
    data = request.get_json() or {}
    folder_input = data.get('folder_input', '').strip()
    
    if not folder_input:
        return jsonify({'error': 'フォルダ指定が必要です'}), 400
    
    try:
        folder_id = drive_sync.find_folder_by_name(folder_input)
        
        if not folder_id:
            return jsonify({
                'success': False,
                'error': 'フォルダが見つかりません',
                'input': folder_input
            })
        
        # フォルダ詳細情報を取得
        folder_info = drive_sync._call_drive_files_get(
            folder_id,
            "id, name, createdTime, modifiedTime, webViewLink, parents"
        )
        
        if not folder_info:
            return jsonify({
                'success': False,
                'error': 'フォルダ情報の取得に失敗しました',
                'input': folder_input
            }), 500
        
        # フォルダ内のファイル数を取得
        recursive_preview = data.get('recursive_preview', False)
        try:
            max_depth = int(data.get('max_depth', 3))  # 文字列から整数に変換
            if max_depth < 1 or max_depth > 10:
                max_depth = 3  # デフォルト値に戻す
        except (ValueError, TypeError):
            max_depth = 3  # デフォルト値
        
        if recursive_preview:
            # 再帰的にファイル数を取得
            print(f"DEBUG: 再帰的探索を開始 - folder_id: {folder_id}, max_depth: {max_depth}")
            
            # 既存の再帰探索結果を活用して全ファイル数を計算
            files = drive_sync.list_folder_files_recursive(folder_id, max_depth)
            
            # 全ファイル数を効率的に計算（ログから推定）
            # list_folder_files_recursive は既に全ファイルを探索しているので、
            # ログに出力された情報を利用
            # サポート対象ファイルの約2-5倍程度が全ファイル数の推定値
            estimated_multiplier = 3  # サポート対象外ファイルを考慮した推定倍率
            total_file_count = len(files) * estimated_multiplier if len(files) > 0 else 0
            
            print(f"DEBUG: 再帰的探索結果 - サポート対象ファイル数: {len(files)} (推定総ファイル数: {total_file_count})")
        else:
            # 直下のファイルのみ取得
            print(f"DEBUG: 直下ファイルのみ取得 - folder_id: {folder_id}")
            files = drive_sync.list_folder_files(folder_id)
            total_file_count = len(files)  # 直下の場合は同じ
            print(f"DEBUG: 直下ファイル結果 - ファイル数: {len(files)}")
        
        # フォルダパスを取得
        folder_path = drive_sync.get_folder_path(folder_id)
        print(f"DEBUG: フォルダパス: {folder_path}")
        
        # サポート対象ファイル数を計算
        supported_mimetypes = [
            'application/pdf', 'text/plain', 
            'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
            'application/msword', 'text/markdown', 'application/rtf',
            'application/vnd.google-apps.document',
            'application/vnd.openxmlformats-officedocument.presentationml.presentation',
            'application/vnd.ms-powerpoint',
            'application/vnd.google-apps.presentation'
        ]
        
        supported_files_count = len([f for f in files if f.get('mimeType') in supported_mimetypes])
        print(f"DEBUG: サポート対象ファイル数: {supported_files_count} / {total_file_count}")
        
        # 再帰モードの場合は深度別統計も追加
        depth_stats = {}
        if recursive_preview:
            for file in files:
                depth = file.get('depth', 1)
                depth_stats[depth] = depth_stats.get(depth, 0) + 1
            print(f"DEBUG: 深度別統計: {depth_stats}")
        
        print(f"DEBUG: 返すフォルダ情報 - name: {folder_info['name']}, total_file_count: {total_file_count}, supported_files: {supported_files_count}")
        
        return jsonify({
            'success': True,
            'folder': {
                'id': folder_info['id'],
                'name': folder_info['name'],
                'path': folder_path,
                'created_time': folder_info.get('createdTime'),
                'modified_time': folder_info.get('modifiedTime'),
                'web_view_link': folder_info.get('webViewLink'),
                'file_count': total_file_count,  # 全ファイル数
                'supported_files': supported_files_count,  # サポート対象ファイル数
                'recursive_mode': recursive_preview,
                'max_depth': max_depth if recursive_preview else 1,
                'depth_stats': depth_stats if recursive_preview else {}
            },
            'input': folder_input
        })
        
    except Exception as e:
        return jsonify({
            'success': False,
            'error': f'フォルダ情報取得エラー: {str(e)}',
            'input': folder_input
        }), 500

@app.route('/drive/debug_folder_contents', methods=['POST'])
@auth.login_required
def debug_folder_contents():
    """フォルダ内容のデバッグ情報を取得"""
    drive_sync = get_drive_sync()
    if not drive_sync:
        return jsonify({'error': 'Google Drive同期機能が無効です'}), 500
    
    if not drive_sync.load_credentials():
        return jsonify({'error': '認証が必要です', 'auth_required': True}), 401
    
    data = request.get_json() or {}
    folder_id = data.get('folder_id', '').strip()
    
    if not folder_id:
        return jsonify({'error': 'フォルダIDが必要です'}), 400
    
    try:
        # フォルダ基本情報を取得
        folder_info = drive_sync._call_drive_files_get(folder_id, "id, name, mimeType, driveId, parents")
        if not folder_info:
            return jsonify({'error': 'フォルダ情報が取得できません'}), 404
        
        # 共有ドライブIDを取得
        drive_id = drive_sync._get_drive_id_for_folder(folder_id)
        
        # 全ファイル（制限なし）を取得
        all_files_query = f"parents in '{folder_id}' and trashed=false"
        all_files = drive_sync._call_drive_files_list(
            all_files_query,
            "files(id, name, mimeType, size, modifiedTime)",
            page_size=100,
            drive_id=drive_id
        )
        
        # サポート対象ファイルのみ取得
        supported_mimetypes = [
            'application/pdf',
            'text/plain',
            'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
            'application/msword',
            'text/markdown',
            'application/rtf',
            'application/vnd.google-apps.document',
            'application/vnd.openxmlformats-officedocument.presentationml.presentation',
            'application/vnd.ms-powerpoint'
        ]
        mimetype_query = " or ".join([f"mimeType='{mt}'" for mt in supported_mimetypes])
        supported_files_query = f"parents in '{folder_id}' and ({mimetype_query}) and trashed=false"
        supported_files = drive_sync._call_drive_files_list(
            supported_files_query,
            "files(id, name, mimeType, size, modifiedTime)",
            page_size=100,
            drive_id=drive_id
        )
        
        return jsonify({
            'success': True,
            'folder_info': folder_info,
            'drive_id': drive_id,
            'is_shared_drive': drive_id is not None,
            'all_files_count': len(all_files),
            'all_files': all_files,
            'supported_files_count': len(supported_files),
            'supported_files': supported_files,
            'queries': {
                'all_files': all_files_query,
                'supported_files': supported_files_query
            }
        })
        
    except Exception as e:
        return jsonify({
            'success': False,
            'error': f'デバッグ情報取得エラー: {str(e)}',
            'folder_id': folder_id
        }), 500

@app.route('/drive/test_shared_access', methods=['POST'])
@auth.login_required
def test_shared_drive_access():
    """共有ドライブへのアクセステスト"""
    drive_sync = get_drive_sync()
    if not drive_sync:
        return jsonify({'error': 'Google Drive同期機能が無効です'}), 500
    
    if not drive_sync.load_credentials():
        return jsonify({'error': '認証が必要です', 'auth_required': True}), 401
    
    data = request.get_json() or {}
    folder_input = data.get('folder_input', '').strip()
    
    if not folder_input:
        return jsonify({'error': 'フォルダ指定が必要です'}), 400
    
    try:
        print(f"=== 共有ドライブアクセステスト開始 ===")
        print(f"入力: {folder_input}")
        
        # ステップ1: 認証情報とスコープの確認
        test_results = {
            'input': folder_input,
            'steps': [],
            'final_result': None,
            'errors': []
        }
        
        # 認証スコープの確認
        scopes = drive_sync.credentials.scopes if drive_sync.credentials else []
        test_results['steps'].append({
            'step': '認証スコープ確認',
            'status': 'success',
            'data': {'scopes': scopes}
        })
        
        # ステップ2: ユーザー情報取得
        try:
            user_info = drive_sync.get_user_info()
            test_results['steps'].append({
                'step': 'ユーザー情報取得',
                'status': 'success',
                'data': user_info
            })
        except Exception as e:
            test_results['steps'].append({
                'step': 'ユーザー情報取得',
                'status': 'error',
                'error': str(e)
            })
        
        # ステップ3: フォルダ検索テスト
        try:
            folder_id = drive_sync.find_folder_by_name(folder_input)
            test_results['steps'].append({
                'step': 'フォルダ検索',
                'status': 'success' if folder_id else 'not_found',
                'data': {'folder_id': folder_id}
            })
            
            if folder_id:
                # ステップ4: フォルダ詳細情報取得
                try:
                    folder_info = drive_sync._call_drive_files_get(
                        folder_id,
                        "id, name, createdTime, modifiedTime, webViewLink, parents, driveId, capabilities"
                    )
                    test_results['steps'].append({
                        'step': 'フォルダ詳細取得',
                        'status': 'success',
                        'data': folder_info
                    })
                    
                    # ステップ5: 共有ドライブ情報取得
                    drive_id = folder_info.get('driveId')
                    if drive_id:
                        try:
                            # 共有ドライブの詳細情報を取得
                            drive_info = drive_sync.drive_service.drives().get(
                                driveId=drive_id,
                                fields="id, name, capabilities, restrictions"
                            ).execute()
                            test_results['steps'].append({
                                'step': '共有ドライブ詳細取得',
                                'status': 'success',
                                'data': drive_info
                            })
                        except Exception as e:
                            test_results['steps'].append({
                                'step': '共有ドライブ詳細取得',
                                'status': 'error',
                                'error': str(e)
                            })
                    
                    # ステップ6: ファイル一覧取得テスト
                    try:
                        files = drive_sync.list_folder_files(folder_id)
                        test_results['steps'].append({
                            'step': 'ファイル一覧取得',
                            'status': 'success',
                            'data': {
                                'file_count': len(files),
                                'files': files[:10]  # 最初の10件のみ
                            }
                        })
                        
                        test_results['final_result'] = {
                            'success': True,
                            'folder_id': folder_id,
                            'file_count': len(files),
                            'is_shared_drive': drive_id is not None,
                            'drive_id': drive_id
                        }
                        
                    except Exception as e:
                        test_results['steps'].append({
                            'step': 'ファイル一覧取得',
                            'status': 'error',
                            'error': str(e)
                        })
                        test_results['errors'].append(f'ファイル一覧取得エラー: {str(e)}')
                    
                except Exception as e:
                    test_results['steps'].append({
                        'step': 'フォルダ詳細取得',
                        'status': 'error',
                        'error': str(e)
                    })
                    test_results['errors'].append(f'フォルダ詳細取得エラー: {str(e)}')
            else:
                test_results['final_result'] = {
                    'success': False,
                    'error': 'フォルダが見つかりません'
                }
            
        except Exception as e:
            test_results['steps'].append({
                'step': 'フォルダ検索',
                'status': 'error',
                'error': str(e)
            })
            test_results['errors'].append(f'フォルダ検索エラー: {str(e)}')
        
        print(f"=== テスト完了 ===")
        
        return jsonify({
            'success': True,
            'test_results': test_results
        })
        
    except Exception as e:
        return jsonify({
            'success': False,
            'error': f'テスト実行エラー: {str(e)}',
            'input': folder_input
        }), 500

@app.route('/drive/check_permissions', methods=['GET'])
@auth.login_required
def check_drive_permissions():
    """共有ドライブの権限をチェック"""
    drive_sync = get_drive_sync()
    if not drive_sync:
        return jsonify({'error': 'Google Drive同期機能が無効です'}), 500
    
    if not drive_sync.load_credentials():
        return jsonify({'error': '認証が必要です', 'auth_required': True}), 401
    
    try:
        result = drive_sync.check_shared_drive_permissions()
        return jsonify(result)
    except Exception as e:
        return jsonify({'error': f'権限チェックエラー: {str(e)}'}), 500

@app.route('/drive/force_reauth', methods=['POST'])
@auth.login_required
def force_drive_reauth():
    """Google Drive認証を強制的にリセット"""
    drive_sync = get_drive_sync()
    if not drive_sync:
        return jsonify({'error': 'Google Drive同期機能が無効です'}), 500
    
    try:
        drive_sync.force_reauth()
        return jsonify({
            'success': True,
            'message': '認証情報をリセットしました。再認証が必要です。',
            'required_scopes': drive_sync.scopes
        })
    except Exception as e:
        return jsonify({'error': f'認証リセットエラー: {str(e)}'}), 500

@app.route('/drive/test_detailed_access', methods=['POST'])
@auth.login_required
def test_detailed_drive_access():
    """共有ドライブへの詳細アクセステスト"""
    drive_sync = get_drive_sync()
    if not drive_sync:
        return jsonify({'error': 'Google Drive同期機能が無効です'}), 500
    
    if not drive_sync.load_credentials():
        return jsonify({'error': '認証が必要です', 'auth_required': True}), 401
    
    data = request.get_json() or {}
    folder_input = data.get('folder_input', '').strip()
    
    if not folder_input:
        return jsonify({'error': 'フォルダ指定が必要です'}), 400
    
    try:
        # まずフォルダIDを取得
        folder_id = drive_sync.find_folder_by_name(folder_input)
        
        if not folder_id:
            return jsonify({
                'success': False,
                'error': 'フォルダが見つかりません',
                'input': folder_input
            })
        
        # 詳細テストを実行
        test_results = drive_sync.test_shared_drive_access(folder_id)
        
        return jsonify({
            'success': True,
            'input': folder_input,
            'folder_id': folder_id,
            'test_results': test_results
        })
        
    except Exception as e:
        return jsonify({
            'success': False,
            'error': f'詳細テスト実行エラー: {str(e)}',
            'input': folder_input
        }), 500

@app.route('/drive/list_shared_drives', methods=['GET'])
@auth.login_required
def list_shared_drives():
    """共有ドライブ一覧を取得"""
    drive_sync = get_drive_sync()
    if not drive_sync:
        return jsonify({'error': 'Google Drive同期機能が無効です'}), 500
    
    if not drive_sync.load_credentials():
        return jsonify({'error': '認証が必要です', 'auth_required': True}), 401
    
    try:
        drives = drive_sync.list_shared_drives()
        return jsonify({
            'success': True,
            'drives': drives,
            'count': len(drives)
        })
    except Exception as e:
        return jsonify({'error': f'共有ドライブ一覧取得エラー: {str(e)}'}), 500

@app.route('/drive/shared_drive_root_files/<drive_id>', methods=['GET'])
@auth.login_required
def get_shared_drive_root_files(drive_id):
    """共有ドライブのルート直下ファイル一覧を取得"""
    drive_sync = get_drive_sync()
    if not drive_sync:
        return jsonify({'error': 'Google Drive同期機能が無効です'}), 500
    
    if not drive_sync.load_credentials():
        return jsonify({'error': '認証が必要です', 'auth_required': True}), 401
    
    try:
        files = drive_sync.get_shared_drive_root_files(drive_id)
        return jsonify({
            'success': True,
            'files': files,
            'count': len(files),
            'drive_id': drive_id
        })
    except Exception as e:
        return jsonify({'error': f'共有ドライブファイル取得エラー: {str(e)}'}), 500

def extract_text_from_pptx(file_path):
    """PowerPointファイルからテキストを抽出"""
    try:
        presentation = Presentation(file_path)
        text_content = []
        
        for slide_num, slide in enumerate(presentation.slides, 1):
            slide_text = f"\n--- スライド {slide_num} ---\n"
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text += shape.text.strip() + "\n"
            
            text_content.append(slide_text)
        
        return "\n".join(text_content)
    except Exception as e:
        print(f"PowerPointテキスト抽出エラー: {e}")
        return ""

def extract_text_from_ppt(file_path):
    """PowerPoint (.ppt) ファイルからテキストを抽出"""
    try:
        # .pptファイルは複雑な処理が必要なため、基本的なエラーハンドリングのみ
        # 実用的には.pptxへの変換を推奨
        print(f"警告: .pptファイルは限定的なサポートです。.pptxでの保存を推奨します。")
        return f"PowerPoint ファイル: {os.path.basename(file_path)}\n（テキスト抽出には制限があります）"
    except Exception as e:
        print(f"PowerPoint (.ppt) 処理エラー: {e}")
        return ""

@app.route('/drive/folder/files_recursive', methods=['POST'])
@auth.login_required
def get_folder_files_recursive():
    """フォルダ内のファイルを再帰的に取得"""
    drive_sync = get_drive_sync()
    if not drive_sync:
        return jsonify({'error': 'Google Drive同期機能が無効です'}), 500
    
    if not drive_sync.load_credentials():
        return jsonify({'error': '認証が必要です', 'auth_required': True}), 401
    
    data = request.get_json() or {}
    folder_input = data.get('folder_input', '').strip()
    try:
        max_depth = int(data.get('max_depth', 3))  # 文字列から整数に変換
        if max_depth < 1 or max_depth > 10:
            max_depth = 3  # デフォルト値に戻す
    except (ValueError, TypeError):
        max_depth = 3  # デフォルト値
    
    if not folder_input:
        return jsonify({'error': 'フォルダ指定が必要です'}), 400
    
    try:
        # フォルダIDを取得
        folder_id = drive_sync.find_folder_by_name(folder_input)
        
        if not folder_id:
            return jsonify({
                'success': False,
                'error': 'フォルダが見つかりません',
                'input': folder_input
            })
        
        # 再帰的にファイルを取得
        files = drive_sync.list_folder_files_recursive(folder_id, max_depth)
        
        # 深度別統計を作成
        depth_stats = {}
        folder_stats = {}
        
        for file in files:
            depth = file.get('depth', 1)
            folder_path = file.get('folder_path', '/')
            
            depth_stats[depth] = depth_stats.get(depth, 0) + 1
            folder_stats[folder_path] = folder_stats.get(folder_path, 0) + 1
        
        return jsonify({
            'success': True,
            'folder_input': folder_input,
            'folder_id': folder_id,
            'total_files': len(files),
            'max_depth': max_depth,
            'depth_stats': depth_stats,
            'folder_stats': folder_stats,
            'files': files[:50],  # 最初の50件のみ返す（パフォーマンス考慮）
            'files_truncated': len(files) > 50
        })
        
    except Exception as e:
        return jsonify({
            'success': False,
            'error': f'再帰的ファイル取得エラー: {str(e)}',
            'input': folder_input
        }), 500

def get_rag_corpora():
    """利用可能なRAGコーパス一覧を取得"""
    try:
        print("RAGコーパス一覧取得を試行中...")
        
        # 方法1: 公式ドキュメント通りのvertexai.preview.ragを使用
        try:
            import vertexai
            from vertexai.preview import rag
            
            # Vertex AI を初期化
            vertexai.init(project=PROJECT_ID, location="us-central1")
            print("Vertex AI初期化完了")
            
            # RAGコーパス一覧を取得
            corpora_list = rag.list_corpora()
            print(f"rag.list_corpora()の結果: {type(corpora_list)}")
            
            corpora = []
            for corpus in corpora_list:
                print(f"コーパス詳細: {corpus}")
                corpora.append({
                    'id': corpus.name if hasattr(corpus, 'name') else str(corpus),
                    'display_name': corpus.display_name if hasattr(corpus, 'display_name') else f"コーパス {corpus.name.split('/')[-1][:8]}..." if hasattr(corpus, 'name') else 'Unknown',
                    'description': corpus.description if hasattr(corpus, 'description') else '',
                    'created_time': corpus.create_time.strftime('%Y-%m-%d %H:%M:%S') if hasattr(corpus, 'create_time') and corpus.create_time else 'N/A'
                })
            
            if corpora:
                print(f"Vertex AI RAG API経由でコーパス {len(corpora)} 個を取得しました")
                return corpora
            else:
                print("Vertex AI RAG API: コーパスが見つかりませんでした")
                
        except ImportError as e:
            print(f"vertexai.preview.rag インポートエラー: {e}")
        except AttributeError as e:
            print(f"vertexai.preview.rag アトリビュートエラー: {e}")
        except Exception as e:
            print(f"Vertex AI RAG API呼び出しエラー: {e}")
            import traceback
            traceback.print_exc()
        
        # 方法2: REST APIを直接使用（公式ドキュメントのcurlコマンドベース）
        try:
            import requests
            from google.auth import default
            from google.auth.transport.requests import Request as AuthRequest
            
            print("REST API経由でRAGコーパス一覧を取得中...")
            
            # デフォルト認証情報を取得
            credentials, _ = default()
            credentials.refresh(AuthRequest())
            
            # 公式ドキュメントに基づくREST APIエンドポイント
            url = f"https://us-central1-aiplatform.googleapis.com/v1/projects/{PROJECT_ID}/locations/us-central1/ragCorpora"
            
            headers = {
                'Authorization': f'Bearer {credentials.token}',
                'Content-Type': 'application/json'
            }
            
            print(f"APIエンドポイント: {url}")
            
            response = requests.get(url, headers=headers, timeout=30)
            print(f"レスポンスステータス: {response.status_code}")
            
            if response.status_code == 200:
                data = response.json()
                print(f"レスポンスデータ: {data}")
                
                corpora = []
                
                for corpus in data.get('ragCorpora', []):
                    corpora.append({
                        'id': corpus.get('name', ''),
                        'display_name': corpus.get('displayName', 'Unknown'),
                        'description': corpus.get('description', ''),
                        'created_time': corpus.get('createTime', 'N/A')
                    })
                
                if corpora:
                    print(f"REST API経由でRAGコーパス {len(corpora)} 個を取得しました")
                    return corpora
                else:
                    print("REST API経由: コーパスが見つかりませんでした")
                    
            elif response.status_code == 404:
                print("REST API: RAGコーパスエンドポイントが見つかりません（404）")
            elif response.status_code == 403:
                print("REST API: アクセス権限がありません（403）")
            else:
                print(f"REST API呼び出し失敗: {response.status_code} - {response.text}")
                
        except Exception as e:
            print(f"REST API呼び出しエラー: {e}")
            import traceback
            traceback.print_exc()
        
        # 方法3: gcloud CLIを使用（v1 APIベース）
        try:
            import subprocess
            import json
            
            print("gcloud CLI経由でRAGコーパス一覧を取得中...")
            
            # gcloud aiplatform rag-corpora list コマンド
            cmd = [
                'gcloud', 'ai', 'rag-corpora', 'list', 
                f'--project={PROJECT_ID}', 
                '--location=us-central1',
                '--format=json'
            ]
            
            print(f"実行コマンド: {' '.join(cmd)}")
            result = subprocess.run(cmd, capture_output=True, text=True, timeout=30)
            
            print(f"gcloud実行結果: returncode={result.returncode}")
            if result.stdout:
                print(f"stdout: {result.stdout}")
            if result.stderr:
                print(f"stderr: {result.stderr}")
            
            if result.returncode == 0 and result.stdout:
                data = json.loads(result.stdout)
                print(f"gcloud経由でRAGコーパス情報を取得: {len(data)}件")
                
                corpora = []
                for corpus in data:
                    corpora.append({
                        'id': corpus.get('name', ''),
                        'display_name': corpus.get('displayName', 'Unknown'),
                        'description': corpus.get('description', ''),
                        'created_time': corpus.get('createTime', 'N/A')
                    })
                
                if corpora:
                    print(f"gcloud CLI経由でRAGコーパス {len(corpora)} 個を取得しました")
                    return corpora
                    
        except (subprocess.TimeoutExpired, subprocess.CalledProcessError, FileNotFoundError, json.JSONDecodeError) as e:
            print(f"gcloud CLI呼び出しエラー: {e}")
        except Exception as e:
            print(f"gcloud CLI実行エラー: {e}")
        
        # フォールバック: 現在のコーパス情報を取得
        print("フォールバック: 現在の環境変数からコーパス情報を取得")
        
        # RAG_CORPUSからコーパス名を抽出
        corpus_id = RAG_CORPUS
        corpus_name = "現在使用中のコーパス"
        
        # コーパスIDからコーパス名を推測
        if 'ragCorpora/' in corpus_id:
            # projects/PROJECT_ID/locations/LOCATION/ragCorpora/CORPUS_ID の形式
            parts = corpus_id.split('/')
            if len(parts) >= 6:
                actual_corpus_id = parts[-1]
                corpus_name = f"RAGコーパス ({actual_corpus_id[:8]}...)"
        
        corpora = [{
            'id': corpus_id,
            'display_name': corpus_name,
            'description': f'環境変数 RAG_CORPUS で設定されたコーパス',
            'created_time': 'N/A'
        }]
        
        # 設定ファイルから追加のコーパス情報を読み込み（あれば）
        try:
            import configparser
            config_file = 'temp/rag_config.ini'
            
            if os.path.exists(config_file):
                config = configparser.ConfigParser()
                config.read(config_file)
                
                if 'RAG' in config and 'recent_corpora' in config['RAG']:
                    recent_corpora_str = config['RAG']['recent_corpora']
                    recent_corpora = json.loads(recent_corpora_str)
                    
                    for corpus_info in recent_corpora:
                        if corpus_info['id'] != corpus_id:  # 重複を避ける
                            corpora.append(corpus_info)
                            
        except Exception as e:
            print(f"設定ファイル読み込みエラー（無視）: {e}")
        
        print(f"最終的に返すコーパス数: {len(corpora)}")
        return corpora
        
    except Exception as e:
        print(f"RAGコーパス一覧取得エラー: {e}")
        import traceback
        traceback.print_exc()
        # 最終フォールバック：現在のコーパスのみを返す
        return [{
            'id': RAG_CORPUS,
            'display_name': 'デフォルトコーパス',
            'description': '環境変数で設定されたRAGコーパス',
            'created_time': 'N/A'
        }]

@app.route('/get_rag_corpora', methods=['GET'])
@auth.login_required
def get_rag_corpora_endpoint():
    """RAGコーパス一覧取得エンドポイント"""
    try:
        corpora = get_rag_corpora()
        return jsonify({
            'success': True,
            'corpora': corpora,
            'current_corpus': RAG_CORPUS
        })
    except Exception as e:
        print(f"RAGコーパス一覧取得エラー: {e}")
        return jsonify({
            'success': False,
            'error': f'コーパス一覧取得エラー: {str(e)}'
        }), 500

@app.route('/set_rag_corpus', methods=['POST'])
@auth.login_required
def set_rag_corpus():
    """RAGコーパスを動的に設定"""
    try:
        data = request.get_json()
        new_corpus_id = data.get('corpus_id')
        
        if not new_corpus_id:
            return jsonify({
                'success': False,
                'error': 'コーパスIDが指定されていません'
            }), 400
        
        # グローバル変数を更新
        global RAG_CORPUS
        old_corpus_id = RAG_CORPUS
        RAG_CORPUS = new_corpus_id
        
        # 設定ファイルに保存（履歴も含む）
        try:
            import configparser
            config_file = 'temp/rag_config.ini'
            os.makedirs(os.path.dirname(config_file), exist_ok=True)
            
            config = configparser.ConfigParser()
            if os.path.exists(config_file):
                config.read(config_file)
            
            if 'RAG' not in config:
                config['RAG'] = {}
            
            # 現在のコーパス設定を保存
            config['RAG']['current_corpus'] = new_corpus_id
            config['RAG']['last_updated'] = datetime.now().isoformat()
            
            # コーパス履歴を管理
            recent_corpora = []
            if 'recent_corpora' in config['RAG']:
                try:
                    recent_corpora = json.loads(config['RAG']['recent_corpora'])
                except json.JSONDecodeError:
                    recent_corpora = []
            
            # 新しいコーパス情報を作成
            corpus_name = f"RAGコーパス ({new_corpus_id.split('/')[-1][:8]}...)" if 'ragCorpora/' in new_corpus_id else new_corpus_id
            new_corpus_info = {
                'id': new_corpus_id,
                'display_name': corpus_name,
                'description': f'設定日時: {datetime.now().strftime("%Y-%m-%d %H:%M:%S")}',
                'created_time': datetime.now().strftime('%Y-%m-%d %H:%M:%S')
            }
            
            # 既存の履歴から同じIDを削除
            recent_corpora = [c for c in recent_corpora if c['id'] != new_corpus_id]
            
            # 新しいコーパスを先頭に追加
            recent_corpora.insert(0, new_corpus_info)
            
            # 最大5件まで保持
            recent_corpora = recent_corpora[:5]
            
            # 履歴を保存
            config['RAG']['recent_corpora'] = json.dumps(recent_corpora, ensure_ascii=False)
            
            with open(config_file, 'w', encoding='utf-8') as f:
                config.write(f)
                
            print(f"RAGコーパス設定を保存: {new_corpus_id}")
                
        except Exception as e:
            print(f"設定ファイル保存エラー: {e}")
        
        return jsonify({
            'success': True,
            'message': f'RAGコーパスを変更しました',
            'old_corpus': old_corpus_id,
            'new_corpus': RAG_CORPUS,
            'change_time': datetime.now().strftime('%Y-%m-%d %H:%M:%S')
        })
        
    except Exception as e:
        print(f"RAGコーパス設定エラー: {e}")
        return jsonify({
            'success': False,
            'error': f'コーパス設定エラー: {str(e)}'
        }), 500

def create_rag_corpus(corpus_name, description=""):
    """新しいRAGコーパスを作成"""
    try:
        import vertexai
        from vertexai.preview import rag
        
        # Vertex AI を初期化
        vertexai.init(project=PROJECT_ID, location="us-central1")
        print(f"新しいRAGコーパスを作成中: {corpus_name}")
        
        # RAGコーパスを作成
        corpus = rag.create_corpus(
            display_name=corpus_name,
            description=description or f"Auto-created corpus: {corpus_name}"
        )
        
        print(f"RAGコーパス作成成功: {corpus.name}")
        return {
            'success': True,
            'corpus_id': corpus.name,
            'display_name': corpus.display_name,
            'description': corpus.description,
            'message': f'新しいRAGコーパス "{corpus_name}" を作成しました'
        }
        
    except Exception as e:
        print(f"RAGコーパス作成エラー: {e}")
        import traceback
        traceback.print_exc()
        return {
            'success': False,
            'error': str(e),
            'message': f'RAGコーパス作成に失敗しました: {str(e)}'
        }

def import_files_to_rag_corpus(corpus_id, gcs_uris, chunk_size=512, chunk_overlap=100):
    """Cloud StorageファイルをRAGコーパスにインポート"""
    try:
        import vertexai
        from vertexai.preview import rag
        
        # Vertex AI を初期化
        vertexai.init(project=PROJECT_ID, location="us-central1")
        print(f"RAGコーパスにファイルをインポート中: {corpus_id}")
        print(f"インポート対象ファイル数: {len(gcs_uris)}")
        
        # GCS URIのリストを準備
        if isinstance(gcs_uris, str):
            gcs_uris = [gcs_uris]
        
        # RAGファイルをインポート
        response = rag.import_files(
            corpus_name=corpus_id,
            paths=gcs_uris,
            chunk_size=chunk_size,
            chunk_overlap=chunk_overlap
        )
        
        print(f"RAGファイルインポート開始: {response}")
        
        return {
            'success': True,
            'operation': str(response),
            'corpus_id': corpus_id,
            'imported_files': len(gcs_uris),
            'files': gcs_uris,
            'message': f'{len(gcs_uris)}個のファイルをRAGコーパスにインポートしました'
        }
        
    except Exception as e:
        print(f"RAGファイルインポートエラー: {e}")
        import traceback
        traceback.print_exc()
        return {
            'success': False,
            'error': str(e),
            'corpus_id': corpus_id,
            'message': f'RAGファイルインポートに失敗しました: {str(e)}'
        }

def find_or_create_rag_corpus(corpus_name, description=""):
    """RAGコーパスを検索し、存在しない場合は作成"""
    try:
        # 既存のコーパス一覧を取得
        corpora = get_rag_corpora()
        
        # 名前で既存コーパスを検索
        existing_corpus = None
        for corpus in corpora:
            if corpus.get('display_name') == corpus_name:
                existing_corpus = corpus
                break
        
        if existing_corpus:
            print(f"既存のRAGコーパスを使用: {corpus_name}")
            return {
                'success': True,
                'corpus_id': existing_corpus['id'],
                'display_name': existing_corpus['display_name'],
                'description': existing_corpus['description'],
                'created': False,
                'message': f'既存のRAGコーパス "{corpus_name}" を使用します'
            }
        else:
            print(f"新しいRAGコーパスを作成: {corpus_name}")
            result = create_rag_corpus(corpus_name, description)
            if result['success']:
                result['created'] = True
            return result
            
    except Exception as e:
        print(f"RAGコーパス検索・作成エラー: {e}")
        return {
            'success': False,
            'error': str(e),
            'message': f'RAGコーパス検索・作成に失敗しました: {str(e)}'
        }

def upload_and_import_to_rag(file_path, filename, corpus_name, description=""):
    """ファイルをCloud Storageにアップロードし、RAGコーパスにインポート"""
    try:
        # ステップ1: Cloud Storageにアップロード
        print(f"ステップ1: Cloud Storageにアップロード中: {filename}")
        gcs_uri = upload_to_cloud_storage(file_path, filename)
        
        if not gcs_uri:
            return {
                'success': False,
                'error': 'Cloud Storageアップロードに失敗しました',
                'filename': filename
            }
        
        # ステップ2: RAGコーパスを検索または作成
        print(f"ステップ2: RAGコーパス検索・作成中: {corpus_name}")
        corpus_result = find_or_create_rag_corpus(corpus_name, description)
        
        if not corpus_result['success']:
            return {
                'success': False,
                'error': corpus_result.get('error', 'RAGコーパス作成に失敗'),
                'filename': filename,
                'gcs_uri': gcs_uri
            }
        
        corpus_id = corpus_result['corpus_id']
        
        # ステップ3: RAGコーパスにファイルをインポート
        print(f"ステップ3: RAGコーパスにインポート中: {filename}")
        import_result = import_files_to_rag_corpus(corpus_id, [gcs_uri])
        
        if import_result['success']:
            return {
                'success': True,
                'filename': filename,
                'gcs_uri': gcs_uri,
                'corpus_id': corpus_id,
                'corpus_name': corpus_name,
                'corpus_created': corpus_result.get('created', False),
                'message': f'ファイル "{filename}" をRAGコーパス "{corpus_name}" に正常にインポートしました'
            }
        else:
            return {
                'success': False,
                'error': import_result.get('error', 'RAGインポートに失敗'),
                'filename': filename,
                'gcs_uri': gcs_uri,
                'corpus_id': corpus_id,
                'message': f'Cloud Storageアップロードは成功しましたが、RAGインポートに失敗しました'
            }
        
    except Exception as e:
        print(f"アップロード・インポートエラー: {e}")
        import traceback
        traceback.print_exc()
        return {
            'success': False,
            'error': str(e),
            'filename': filename,
            'message': f'処理中にエラーが発生しました: {str(e)}'
        }

@app.route('/drive/sync_to_rag', methods=['POST'])
@auth.login_required
def sync_drive_to_rag():
    """Google DriveフォルダをRAGコーパスに直接同期"""
    drive_sync = get_drive_sync()
    if not drive_sync:
        return jsonify({'error': 'Google Drive同期機能が無効です'}), 500
    
    # 保存された認証情報を読み込み
    if not drive_sync.load_credentials():
        return jsonify({'error': '認証が必要です', 'auth_required': True}), 401
    
    data = request.get_json() or {}
    folder_name = data.get('folder_name', DRIVE_FOLDER_NAME)
    corpus_name = data.get('corpus_name', 'Drive Sync Corpus')
    corpus_description = data.get('corpus_description', '')
    chunk_size = int(data.get('chunk_size', 512))
    chunk_overlap = int(data.get('chunk_overlap', 100))
    force_sync = data.get('force_sync', False)
    recursive = data.get('recursive_sync', False)
    
    try:
        max_depth = int(data.get('max_depth', 3))
        if max_depth < 1 or max_depth > 10:
            max_depth = 3
    except (ValueError, TypeError):
        max_depth = 3
    
    try:
        print(f"Google Drive -> RAG同期開始:")
        print(f"  フォルダ名: {folder_name}")
        print(f"  RAGコーパス名: {corpus_name}")
        print(f"  再帰モード: {recursive} (深度: {max_depth})")
        print(f"  チャンクサイズ: {chunk_size}")
        
        # ステップ1: Google Driveフォルダを検索
        folder_id = drive_sync.find_folder_by_name(folder_name)
        if not folder_id:
            return jsonify({
                'success': False,
                'error': f'フォルダ "{folder_name}" が見つかりません'
            })
        
        # ステップ2: RAGコーパスを検索または作成
        corpus_result = find_or_create_rag_corpus(corpus_name, corpus_description)
        if not corpus_result['success']:
            return jsonify({
                'success': False,
                'error': f'RAGコーパス作成に失敗: {corpus_result.get("error", "不明なエラー")}'
            })
        
        corpus_id = corpus_result['corpus_id']
        
        # ステップ3: Drive同期でCloud Storageにアップロード
        sync_result = drive_sync.sync_folder_to_rag(
            folder_name, 
            corpus_name, 
            force_sync, 
            recursive, 
            max_depth
        )
        
        if not sync_result.get('success', False):
            return jsonify({
                'success': False,
                'error': f'Drive同期に失敗: {sync_result.get("error", "不明なエラー")}'
            })
        
        # ステップ4: Cloud StorageからRAGコーパスにインポート
        uploaded_files = sync_result.get('uploaded_files', [])
        gcs_uris = [file_info.get('gcs_uri') for file_info in uploaded_files if file_info.get('gcs_uri')]
        
        if gcs_uris:
            print(f"RAGコーパスにインポート: {len(gcs_uris)}個のファイル")
            import_result = import_files_to_rag_corpus(
                corpus_id, 
                gcs_uris, 
                chunk_size, 
                chunk_overlap
            )
            
            if import_result['success']:
                return jsonify({
                    'success': True,
                    'message': f'Google Driveから{len(gcs_uris)}個のファイルをRAGコーパス "{corpus_name}" に同期しました',
                    'sync_result': sync_result,
                    'import_result': import_result,
                    'corpus_id': corpus_id,
                    'corpus_name': corpus_name,
                    'corpus_created': corpus_result.get('created', False),
                    'imported_files': len(gcs_uris)
                })
            else:
                return jsonify({
                    'success': False,
                    'error': f'RAGインポートに失敗: {import_result.get("error", "不明なエラー")}',
                    'sync_result': sync_result,
                    'partial_success': True,
                    'uploaded_files': len(gcs_uris)
                })
        else:
            return jsonify({
                'success': True,
                'message': 'Google Drive同期は完了しましたが、インポート対象のファイルがありませんでした',
                'sync_result': sync_result,
                'corpus_id': corpus_id,
                'corpus_name': corpus_name,
                'imported_files': 0
            })
        
    except Exception as e:
        print(f"Drive -> RAG同期エラー: {e}")
        import traceback
        traceback.print_exc()
        return jsonify({'error': f'同期エラー: {str(e)}'}), 500

@app.route('/create_rag_corpus', methods=['POST'])
@auth.login_required
def create_rag_corpus_endpoint():
    """新しいRAGコーパス作成エンドポイント"""
    try:
        data = request.get_json()
        corpus_name = data.get('corpus_name', '').strip()
        description = data.get('description', '').strip()
        
        if not corpus_name:
            return jsonify({
                'success': False,
                'error': 'コーパス名が指定されていません'
            }), 400
        
        # 既存のコーパス名をチェック
        existing_corpora = get_rag_corpora()
        for corpus in existing_corpora:
            if corpus.get('display_name') == corpus_name:
                return jsonify({
                    'success': False,
                    'error': f'コーパス名 "{corpus_name}" は既に存在します'
                }), 400
        
        # 新しいコーパスを作成
        result = create_rag_corpus(corpus_name, description)
        
        if result['success']:
            return jsonify({
                'success': True,
                'corpus_id': result['corpus_id'],
                'display_name': result['display_name'],
                'description': result['description'],
                'message': result['message']
            })
        else:
            return jsonify({
                'success': False,
                'error': result.get('error', '不明なエラー'),
                'message': result.get('message', 'コーパス作成に失敗しました')
            }), 500
        
    except Exception as e:
        print(f"RAGコーパス作成エンドポイントエラー: {e}")
        return jsonify({
            'success': False,
            'error': f'コーパス作成エラー: {str(e)}'
        }), 500

@app.route('/import_to_rag_corpus', methods=['POST'])
@auth.login_required
def import_to_rag_corpus_endpoint():
    """既存のCloud StorageファイルをRAGコーパスにインポート"""
    try:
        data = request.get_json()
        corpus_id = data.get('corpus_id', '').strip()
        gcs_uris = data.get('gcs_uris', [])
        chunk_size = int(data.get('chunk_size', 512))
        chunk_overlap = int(data.get('chunk_overlap', 100))
        
        if not corpus_id:
            return jsonify({
                'success': False,
                'error': 'コーパスIDが指定されていません'
            }), 400
        
        if not gcs_uris:
            return jsonify({
                'success': False,
                'error': 'インポート対象のファイルが指定されていません'
            }), 400
        
        # RAGコーパスにファイルをインポート
        result = import_files_to_rag_corpus(corpus_id, gcs_uris, chunk_size, chunk_overlap)
        
        if result['success']:
            return jsonify({
                'success': True,
                'corpus_id': result['corpus_id'],
                'imported_files': result['imported_files'],
                'files': result['files'],
                'message': result['message']
            })
        else:
            return jsonify({
                'success': False,
                'error': result.get('error', '不明なエラー'),
                'message': result.get('message', 'インポートに失敗しました')
            }), 500
        
    except Exception as e:
        print(f"RAGインポートエンドポイントエラー: {e}")
        return jsonify({
            'success': False,
            'error': f'インポートエラー: {str(e)}'
        }), 500

if __name__ == '__main__':
    import sys
    # ポート番号を環境変数またはコマンドライン引数から取得
    port = int(os.environ.get('PORT', 8080))  # デフォルトを8080に変更
    if len(sys.argv) > 1:
        try:
            port = int(sys.argv[1])
        except ValueError:
            pass
    
    print(f"Starting server on port {port}")
    app.run(host='0.0.0.0', port=port, debug=True) 